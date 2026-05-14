"""
Gera apresentação .pptx — Oportunidade de Investimento Imobiliário
Praça da República, 162/166 — Centro Histórico de SP
13 slides, foco em VENDA + Programa Requalifica Centro
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# --- Paths ---
ROOT = Path(__file__).parent
IMG = Path(r"C:\Users\xandao\Desktop\predio venda")
LOGO = Path(r"C:\Users\xandao\Downloads\capa.png")
OUT = ROOT / "Oportunidade_Investimento_Praca_Republica.pptx"

# --- Paleta ---
GRAFITE = RGBColor(0x1F, 0x29, 0x37)
DOURADO = RGBColor(0xC9, 0xA9, 0x61)
DOURADO_ESC = RGBColor(0x8B, 0x6F, 0x3E)
OFFWHITE = RGBColor(0xF5, 0xF1, 0xEA)
CINZA = RGBColor(0xE5, 0xE5, 0xE5)
CINZA_M = RGBColor(0x6B, 0x72, 0x80)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
VERDE = RGBColor(0x16, 0x82, 0x4F)

# Fontes
T_SERIF = "Georgia"     # títulos (vibe patrimônio)
T_SANS = "Calibri"      # corpo

# --- Setup deck 16:9 ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh

def add_text(slide, x, y, w, h, text, *, font=T_SANS, size=14, color=GRAFITE,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_image(slide, path, x, y, w, h, *, crop_to_fill=True):
    """Adiciona imagem preenchendo o frame (crop centralizado)."""
    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if crop_to_fill:
        # python-pptx ajusta forçando width+height — pode distorcer. 
        # Para crop real precisaríamos calcular ratios. Simplificação: aceitar leve distorção.
        pass
    return pic

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def gold_line(slide, x, y, w, thickness_pt=2, color=DOURADO):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = Pt(thickness_pt)
    return line

def page_number(slide, num, total=13, color=DOURADO):
    add_text(slide, Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
             f"{num:02d} / {total:02d}", size=9, color=color, align=PP_ALIGN.RIGHT,
             font=T_SANS)

def footer_brand(slide, color=CINZA_M):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "PRAÇA DA REPÚBLICA · 162/166 · CENTRO HISTÓRICO · SÃO PAULO/SP",
             size=8, color=color, font=T_SANS)

# ============================================================
# SLIDE 1 — CAPA
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, GRAFITE)
# Foto principal ocupando metade direita
add_image(s, IMG / "vista do predio.jpg", Inches(5.5), 0, Inches(7.833), SH)
# Overlay grafite sobre a foto (gradiente simulado com retângulo translúcido)
ov = add_rect(s, Inches(5.5), 0, Inches(2.5), SH, GRAFITE)
ov.fill.transparency = 0  # python-pptx limitação, usar alpha via xml
# Forçar transparência via XML
sp = ov.fill._xPr.find(qn("a:solidFill"))
if sp is not None:
    clr = sp.find(qn("a:srgbClr"))
    if clr is not None:
        alpha = etree.SubElement(clr, qn("a:alpha"))
        alpha.set("val", "70000")  # 70% opaco

# Linha dourada vertical decorativa
add_rect(s, Inches(0.7), Inches(2.7), Inches(0.04), Inches(2.2), DOURADO)

# Textos da capa
add_text(s, Inches(0.9), Inches(2.5), Inches(5), Inches(0.5),
         "DOSSIÊ DE INVESTIMENTO", size=12, color=DOURADO,
         font=T_SANS, bold=True)
add_text(s, Inches(0.9), Inches(2.9), Inches(5), Inches(2.2),
         "Oportunidade\nImobiliária\nem São Paulo",
         font=T_SERIF, size=44, color=OFFWHITE, bold=True, line_spacing=1.05)
add_text(s, Inches(0.9), Inches(5.3), Inches(5), Inches(0.5),
         "Praça da República · 162/166", size=16, color=OFFWHITE,
         font=T_SERIF, italic=True)
add_text(s, Inches(0.9), Inches(5.7), Inches(5), Inches(0.4),
         "Centro Histórico · São Paulo/SP", size=12, color=DOURADO,
         font=T_SANS)

# Rodapé capa
add_text(s, Inches(0.9), Inches(6.9), Inches(5), Inches(0.3),
         "EDIFÍCIO TOMBADO · 4.000 m² · ZONA ZEU · PROGRAMA REQUALIFICA CENTRO",
         size=8.5, color=DOURADO, font=T_SANS, bold=True)

# ============================================================
# SLIDE 2 — SUMÁRIO EXECUTIVO
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "SUMÁRIO EXECUTIVO", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(10), Inches(0.8),
         "Um ativo singular no coração de São Paulo",
         font=T_SERIF, size=28, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

# Texto introdutório
add_text(s, Inches(0.7), Inches(2.1), Inches(6.3), Inches(2.5),
         "Edifício comercial de 4.000 m² localizado na Praça da "
         "República, em frente à estação de metrô (Linhas 3-Vermelha "
         "e 4-Amarela). Antiga sede da Secretaria de Educação do "
         "Estado de São Paulo, é um imóvel tombado, com documentação "
         "consolidada e elegibilidade plena ao Programa Requalifica "
         "Centro — que oferece a combinação mais agressiva de "
         "incentivos fiscais já praticada no município.",
         size=12, color=GRAFITE, line_spacing=1.4)

# Cards de KPI à direita
def kpi_card(x, y, label, value, sub=""):
    add_rect(s, x, y, Inches(2.7), Inches(1.4), BRANCO)
    add_rect(s, x, y, Inches(0.08), Inches(1.4), DOURADO)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(2.5), Inches(0.3),
             label, size=8.5, color=CINZA_M, bold=True, font=T_SANS)
    add_text(s, x + Inches(0.2), y + Inches(0.45), Inches(2.5), Inches(0.6),
             value, size=20, color=GRAFITE, bold=True, font=T_SERIF)
    if sub:
        add_text(s, x + Inches(0.2), y + Inches(1.05), Inches(2.5), Inches(0.3),
                 sub, size=8.5, color=DOURADO_ESC, font=T_SANS, italic=True)

kpi_card(Inches(7.3), Inches(2.1), "ÁREA CONSTRUÍDA", "4.000 m²", "13 andares + sobreloja + mezanino")
kpi_card(Inches(10.2), Inches(2.1), "AVALIAÇÃO", "R$ 31–56 mi", "venda direta")
kpi_card(Inches(7.3), Inches(3.6), "FRAÇÕES", "38 matrículas", "desmembradas")
kpi_card(Inches(10.2), Inches(3.6), "ZONEAMENTO", "ZEU", "uso misto, alta densidade")

# Faixa inferior — destaques
add_rect(s, 0, Inches(5.4), SW, Inches(1.5), GRAFITE)
add_text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.4),
         "PILARES DA OPORTUNIDADE", size=10, color=DOURADO, bold=True, font=T_SANS)

pilares = [
    ("Localização", "Frente à Praça da\nRepública · metrô 100m"),
    ("Patrimônio", "Edifício tombado ·\nselo histórico"),
    ("Estrutura", "38 frações ·\nflexibilidade total"),
    ("Incentivos", "Requalifica Centro ·\nIPTU + ISS + ITBI + 25%"),
]
for i, (h, d) in enumerate(pilares):
    x = Inches(0.7 + i * 3.1)
    add_text(s, x, Inches(5.95), Inches(3), Inches(0.4),
             h.upper(), size=11, color=DOURADO, bold=True, font=T_SANS)
    add_text(s, x, Inches(6.3), Inches(3), Inches(0.8),
             d, size=10, color=OFFWHITE, font=T_SANS, line_spacing=1.3)

page_number(s, 2)
footer_brand(s)

# ============================================================
# SLIDE 3 — O ATIVO (ficha técnica)
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
# Foto à esquerda
add_image(s, IMG / "recepcao + elevadores.jpg", 0, 0, Inches(5.5), SH)
# Painel à direita
add_rect(s, Inches(5.5), 0, Inches(7.833), SH, OFFWHITE)
add_text(s, Inches(5.9), Inches(0.6), Inches(7), Inches(0.4),
         "O ATIVO", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(5.9), Inches(1.0), Inches(7), Inches(0.8),
         "Ficha Técnica", font=T_SERIF, size=30, color=GRAFITE, bold=True)
gold_line(s, Inches(5.9), Inches(1.85), Inches(2), 2.5)

ficha = [
    ("Endereço", "Praça da República, 162/166 — Centro, São Paulo/SP"),
    ("Área Construída", "4.000 m²"),
    ("Estrutura", "13 andares + sobreloja + mezanino"),
    ("Salas", "30 unidades amplas, layout flexível"),
    ("Elevadores", "2 unidades funcionais"),
    ("Capacidade", "Até 600 pessoas simultâneas"),
    ("Zoneamento", "ZEU — Zona de Centralidade"),
    ("Matrículas", "38 frações desmembradas"),
    ("AVCB", "Regular e atualizado"),
    ("Histórico", "Antiga sede da Secretaria de Educação do Estado de SP"),
    ("Tombamento", "Imóvel tombado — patrimônio histórico"),
    ("IPTU", "≈ R$ 13.000 / mês no conjunto"),
]
y = Inches(2.15)
for label, value in ficha:
    add_text(s, Inches(5.9), y, Inches(2.0), Inches(0.32),
             label.upper(), size=9, color=DOURADO_ESC, bold=True, font=T_SANS)
    add_text(s, Inches(7.95), y, Inches(5.2), Inches(0.32),
             value, size=10.5, color=GRAFITE, font=T_SANS)
    y += Inches(0.34)

page_number(s, 3, color=GRAFITE)

# ============================================================
# SLIDE 4 — LOCALIZAÇÃO
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, GRAFITE)
# Foto vista da praça com placa do metrô
add_image(s, IMG / "visao varanda.jpg", Inches(6.5), 0, Inches(6.833), SH)
# Overlay
ov = add_rect(s, Inches(6.5), 0, Inches(1.5), SH, GRAFITE)
sp = ov.fill._xPr.find(qn("a:solidFill"))
clr = sp.find(qn("a:srgbClr"))
alpha = etree.SubElement(clr, qn("a:alpha")); alpha.set("val", "60000")

add_text(s, Inches(0.7), Inches(0.6), Inches(6), Inches(0.4),
         "LOCALIZAÇÃO ESTRATÉGICA", size=11, color=DOURADO, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(1.0), Inches(6), Inches(1.6),
         "No centro do\ncentro de SP",
         font=T_SERIF, size=36, color=OFFWHITE, bold=True, line_spacing=1.0)
gold_line(s, Inches(0.7), Inches(2.85), Inches(2), 2.5)

# Bullets de localização
locs = [
    ("100 m", "da Estação República — Linhas 3-Vermelha e 4-Amarela do Metrô"),
    ("Direto", "para Avenidas Ipiranga, São João, São Luís e Consolação"),
    ("Entorno", "Edifício Copan, Biblioteca Mário de Andrade, Teatro Municipal"),
    ("Perímetro", "de revitalização urbana acelerada (Requalifica Centro)"),
    ("Hub", "comercial, jurídico, educacional, cultural e turístico"),
]
y = Inches(3.2)
for n, d in locs:
    add_text(s, Inches(0.7), y, Inches(1.3), Inches(0.5),
             n, size=18, color=DOURADO, bold=True, font=T_SERIF)
    add_text(s, Inches(2.1), y + Inches(0.05), Inches(4.4), Inches(0.6),
             d, size=11, color=OFFWHITE, font=T_SANS, line_spacing=1.3)
    y += Inches(0.7)

page_number(s, 4)
footer_brand(s, OFFWHITE)

# ============================================================
# SLIDE 5 — LAYOUT & FLEXIBILIDADE
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "LAYOUT & FLEXIBILIDADE", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "30 salas amplas · 38 frações desmembradas",
         font=T_SERIF, size=28, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

add_text(s, Inches(0.7), Inches(2.1), Inches(5.5), Inches(2.2),
         "Cada pavimento foi concebido para receber operações "
         "corporativas de grande porte. As 38 matrículas autônomas "
         "permitem alienação fracionada, garantias segregadas e "
         "estruturação financeira flexível — desde venda integral "
         "para single-tenant até montagem de FII com múltiplos "
         "investidores.",
         size=12, color=GRAFITE, line_spacing=1.45)

# Atributos
attrs = [
    "Plantas livres com colunas espaçadas",
    "Pé-direito generoso",
    "Janelões com luz natural",
    "Iluminação setorizada por andar",
    "Layout adaptável a open space, salas privativas ou auditório",
]
y = Inches(4.4)
for a in attrs:
    add_text(s, Inches(0.85), y, Inches(0.3), Inches(0.3),
             "▎", size=14, color=DOURADO, bold=True, font=T_SANS)
    add_text(s, Inches(1.1), y, Inches(5.2), Inches(0.3),
             a, size=10.5, color=GRAFITE, font=T_SANS)
    y += Inches(0.36)

# Duas imagens à direita
add_image(s, IMG / "sala grande.jpg", Inches(6.7), Inches(2.1), Inches(6.3), Inches(2.5))
add_image(s, IMG / "visao interna + externa 3 ampliada.jpg",
          Inches(6.7), Inches(4.75), Inches(6.3), Inches(2.5))

page_number(s, 5, color=GRAFITE)

# ============================================================
# SLIDE 6 — INFRAESTRUTURA
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "INFRAESTRUTURA", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "Pronto para operar",
         font=T_SERIF, size=28, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

# 3 imagens em fila
add_image(s, IMG / "elevadores.jpg", Inches(0.7), Inches(2.1), Inches(4), Inches(2.6))
add_image(s, IMG / "quadro de luz.jpg", Inches(4.85), Inches(2.1), Inches(4), Inches(2.6))
add_image(s, IMG / "corredores.jpg", Inches(9.0), Inches(2.1), Inches(3.5), Inches(2.6))

# Cards de infra
infos = [
    ("ELEVADORES", "2 unidades funcionais"),
    ("ELÉTRICA", "Quadro setorizado andar a andar"),
    ("CIRCULAÇÃO", "Corredores amplos e escadas duplas"),
    ("AVCB", "Regular e atualizado"),
    ("CLIMATIZAÇÃO", "Split + cassete em uso aprovado"),
    ("ACERVO TÉCNICO", "Plantas e projetos disponíveis"),
]
y_start = Inches(5.0)
for i, (h, d) in enumerate(infos):
    col, row = i % 3, i // 3
    x = Inches(0.7 + col * 4.15)
    y = y_start + Inches(row * 1.1)
    add_rect(s, x, y, Inches(0.06), Inches(0.9), DOURADO)
    add_text(s, x + Inches(0.2), y, Inches(3.8), Inches(0.35),
             h, size=10, color=DOURADO_ESC, bold=True, font=T_SANS)
    add_text(s, x + Inches(0.2), y + Inches(0.35), Inches(3.8), Inches(0.55),
             d, size=10.5, color=GRAFITE, font=T_SANS, line_spacing=1.25)

page_number(s, 6, color=GRAFITE)

# ============================================================
# SLIDE 7 — REQUALIFICA CENTRO (CAPA DO BLOCO)
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, GRAFITE)
add_image(s, IMG / "sobreloja vista rua.jpg", Inches(6.5), 0, Inches(6.833), SH)
ov = add_rect(s, Inches(6.5), 0, Inches(1.0), SH, GRAFITE)
sp = ov.fill._xPr.find(qn("a:solidFill"))
clr = sp.find(qn("a:srgbClr"))
alpha = etree.SubElement(clr, qn("a:alpha")); alpha.set("val", "55000")

add_text(s, Inches(0.7), Inches(0.6), Inches(7), Inches(0.4),
         "A TESE DE INVESTIMENTO", size=11, color=DOURADO, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(1.0), Inches(6.5), Inches(2.4),
         "Programa\nRequalifica\nCentro",
         font=T_SERIF, size=44, color=OFFWHITE, bold=True, line_spacing=0.95)
gold_line(s, Inches(0.7), Inches(4.0), Inches(2), 2.5)

add_text(s, Inches(0.7), Inches(4.25), Inches(6), Inches(2.2),
         "O imóvel está plenamente inserido no perímetro de "
         "benefícios do programa municipal mais agressivo já lançado "
         "para o Centro de São Paulo. A combinação de incentivos "
         "tributários, subvenção direta e licenciamento expresso "
         "transforma a equação financeira do investimento.",
         size=12.5, color=OFFWHITE, line_spacing=1.45, italic=True)

# Base legal pequena
add_text(s, Inches(0.7), Inches(6.5), Inches(6), Inches(0.5),
         "Base legal: Lei Municipal 17.577/2021 · Lei 17.844/2022 (AIU-SCE) · "
         "Decretos 61.311/2022 e 62.465/2023 · IN SF/SUREM nº 01/2023",
         size=8.5, color=DOURADO, font=T_SANS, italic=True)

page_number(s, 7)

# ============================================================
# SLIDE 8 — STACK DE INCENTIVOS FISCAIS
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "REQUALIFICA CENTRO · 1/2", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "Stack de incentivos tributários",
         font=T_SERIF, size=26, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.8), Inches(2), 2.5)

# Tabela de benefícios — 3 colunas grandes (cards)
beneficios = [
    {
        "tag": "IPTU",
        "title": "Imposto Predial",
        "body": [
            "Remissão integral de créditos anteriores",
            "Isenção total por 3 anos após CCR",
            "Alíquotas progressivas do 4º ao 8º ano",
            "Retorno integral apenas no 9º ano",
        ],
        "impacto": "Economia > R$ 2.000.000\nno ciclo de 8 anos",
    },
    {
        "tag": "ISS",
        "title": "Imposto sobre Serviços",
        "body": [
            "Redução de 5% para 2% na alíquota",
            "Aplicável a construção, engenharia,",
            "arquitetura e reforma vinculadas",
            "ao retrofit",
        ],
        "impacto": "Economia de R$ 300.000\nem obra de R$ 10 milhões",
    },
    {
        "tag": "ITBI",
        "title": "Transmissão de Imóveis",
        "body": [
            "Isenção total nas transmissões",
            "destinadas à requalificação",
            "Redução drástica do custo",
            "de aquisição e estruturação",
        ],
        "impacto": "Isenção 100%\nna aquisição",
    },
]

card_w = Inches(4.0)
card_h = Inches(4.5)
gap = Inches(0.25)
total_w = card_w * 3 + gap * 2
start_x = (SW - total_w) / 2
y = Inches(2.1)

for i, b in enumerate(beneficios):
    x = start_x + (card_w + gap) * i
    # Card branco
    add_rect(s, x, y, card_w, card_h, BRANCO)
    # Faixa dourada topo
    add_rect(s, x, y, card_w, Inches(0.5), DOURADO)
    add_text(s, x, y + Inches(0.08), card_w, Inches(0.4),
             b["tag"], size=14, color=BRANCO, bold=True, font=T_SERIF, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + Inches(0.7), card_w - Inches(0.6), Inches(0.5),
             b["title"], size=14, color=GRAFITE, bold=True, font=T_SERIF)
    # Bullets
    yy = y + Inches(1.3)
    for line in b["body"]:
        add_text(s, x + Inches(0.4), yy, Inches(0.2), Inches(0.3),
                 "•", size=12, color=DOURADO, bold=True, font=T_SANS)
        add_text(s, x + Inches(0.6), yy, card_w - Inches(0.8), Inches(0.3),
                 line, size=10.5, color=GRAFITE, font=T_SANS)
        yy += Inches(0.35)
    # Impacto
    add_rect(s, x + Inches(0.3), y + card_h - Inches(1.1),
             card_w - Inches(0.6), Inches(0.85), OFFWHITE)
    add_text(s, x + Inches(0.3), y + card_h - Inches(1.05),
             card_w - Inches(0.6), Inches(0.3),
             "IMPACTO DIRETO", size=8, color=DOURADO_ESC, bold=True, font=T_SANS,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + card_h - Inches(0.75),
             card_w - Inches(0.6), Inches(0.55),
             b["impacto"], size=11, color=GRAFITE, bold=True, font=T_SERIF,
             align=PP_ALIGN.CENTER, line_spacing=1.15)

# Faixa de impacto consolidado
add_rect(s, 0, Inches(6.8), SW, Inches(0.7), GRAFITE)
add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
         "IPTU atual ≈ R$ 13.000/mês · Isenção de 3 anos = R$ 468.000 economizados só em IPTU",
         size=12, color=DOURADO, bold=True, font=T_SANS, align=PP_ALIGN.CENTER)

page_number(s, 8, color=GRAFITE)

# ============================================================
# SLIDE 9 — SUBVENÇÃO DIRETA + LICENCIAMENTO
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "REQUALIFICA CENTRO · 2/2", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "Subvenção direta + via expressa",
         font=T_SERIF, size=26, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.8), Inches(2), 2.5)

# Bloco grande à esquerda: Subvenção 25%
add_rect(s, Inches(0.7), Inches(2.15), Inches(6.3), Inches(4.0), GRAFITE)
add_text(s, Inches(1.0), Inches(2.3), Inches(6), Inches(0.4),
         "SUBVENÇÃO ECONÔMICA", size=10, color=DOURADO, bold=True, font=T_SANS)
add_text(s, Inches(1.0), Inches(2.7), Inches(6), Inches(1.3),
         "Até 25%\na fundo perdido",
         font=T_SERIF, size=34, color=OFFWHITE, bold=True, line_spacing=1.0)
add_text(s, Inches(1.0), Inches(4.5), Inches(6), Inches(1.5),
         "Aporte direto do município sobre o valor das obras, "
         "mediante participação em chamamentos públicos da "
         "Secretaria Municipal de Urbanismo e Licenciamento (SMUL).",
         size=11, color=OFFWHITE, font=T_SANS, line_spacing=1.4)
add_rect(s, Inches(1.0), Inches(5.5), Inches(5.7), Inches(0.5), DOURADO)
add_text(s, Inches(1.0), Inches(5.55), Inches(5.7), Inches(0.4),
         "Em obra de R$ 10 mi → até R$ 2.500.000 a fundo perdido",
         size=11.5, color=GRAFITE, bold=True, font=T_SANS, align=PP_ALIGN.CENTER)

# Bloco direito: Licenciamento + benefícios edilícios
add_text(s, Inches(7.3), Inches(2.3), Inches(6), Inches(0.4),
         "LICENCIAMENTO EXPRESSO", size=10, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(7.3), Inches(2.65), Inches(6), Inches(0.6),
         "Rito sumário · 60 dias úteis",
         font=T_SERIF, size=20, color=GRAFITE, bold=True)

itens = [
    "Aprovação de projetos em até 60 dias úteis",
    "Isenção de taxas municipais de licenciamento por 5 anos",
    "Isenção de taxas de fiscalização no mesmo período",
    "Prioridade em contratos com órgãos públicos",
]
y = Inches(3.5)
for it in itens:
    add_text(s, Inches(7.3), y, Inches(0.3), Inches(0.3),
             "▎", size=14, color=DOURADO, bold=True, font=T_SANS)
    add_text(s, Inches(7.55), y, Inches(5.6), Inches(0.5),
             it, size=11, color=GRAFITE, font=T_SANS, line_spacing=1.3)
    y += Inches(0.5)

# Elegibilidade já garantida
add_rect(s, Inches(7.3), Inches(5.7), Inches(5.7), Inches(1.2), VERDE)
add_text(s, Inches(7.5), Inches(5.8), Inches(5.4), Inches(0.4),
         "✓ ELEGIBILIDADE GARANTIDA", size=10, color=BRANCO, bold=True, font=T_SANS)
add_text(s, Inches(7.5), Inches(6.15), Inches(5.4), Inches(0.7),
         "Edificação anterior a 1992 + Zona ZEU\nArt. 37 da Lei 17.844/2022",
         size=10.5, color=BRANCO, font=T_SANS, line_spacing=1.3)

page_number(s, 9, color=GRAFITE)

# ============================================================
# SLIDE 10 — VOCAÇÕES DE USO
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "VOCAÇÕES DE USO", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "Múltiplos modelos de exploração",
         font=T_SERIF, size=28, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

vocacoes = [
    ("Sede Corporativa", "Para empresas que buscam representatividade, centralidade e endereço de prestígio em São Paulo."),
    ("Fundo Imobiliário (FII)", "Ativo com perfil ideal para composição de portfólio de renda — 38 frações permitem entrada faseada."),
    ("Operação Educacional", "Estrutura compatível com faculdades, centros de treinamento, escolas técnicas e cursos livres."),
    ("Retrofit Mixed-Use", "Conversão para uso residencial, hoteleiro ou misto com incentivos municipais ampliados."),
    ("Family Office", "Posicionamento patrimonial em ativo histórico, com selo cultural e alta liquidez fracionada."),
    ("Sede Institucional", "Vocação natural para órgãos públicos, instituições culturais, fundações ou OSCIPs."),
]
# 2 colunas x 3 linhas
for i, (h, d) in enumerate(vocacoes):
    col, row = i % 2, i // 2
    x = Inches(0.7 + col * 6.3)
    y = Inches(2.15 + row * 1.55)
    add_rect(s, x, y, Inches(6), Inches(1.35), BRANCO)
    add_rect(s, x, y, Inches(0.08), Inches(1.35), DOURADO)
    add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(5.6), Inches(0.45),
             h, size=14, color=GRAFITE, bold=True, font=T_SERIF)
    add_text(s, x + Inches(0.25), y + Inches(0.6), Inches(5.6), Inches(0.75),
             d, size=10.5, color=CINZA_M, font=T_SANS, line_spacing=1.35)

page_number(s, 10, color=GRAFITE)

# ============================================================
# SLIDE 11 — CENÁRIOS DE RETORNO
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "RENTABILIDADE", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "Cenários projetados de retorno",
         font=T_SERIF, size=28, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

add_text(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.5),
         "Simulação pós-aquisição e retrofit · valorização do m² na região central",
         size=11, color=CINZA_M, italic=True, font=T_SANS)

# Tabela visual — 4 cards (cenário)
cenarios = [
    ("Conservador", "R$ 31 mi",  "5,8%", "58%",  CINZA),
    ("Moderado",    "R$ 43,5 mi","6,1%", "61%",  DOURADO),
    ("Realista",    "R$ 43,5 mi","7,4%", "74%",  DOURADO_ESC),
    ("Otimista",    "R$ 56 mi",  "6,9%", "69%",  GRAFITE),
]
card_w = Inches(2.9)
card_h = Inches(3.6)
gap = Inches(0.2)
start_x = (SW - (card_w * 4 + gap * 3)) / 2
y = Inches(2.8)

for i, (nome, preco, cap, roi, accent) in enumerate(cenarios):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, y, card_w, card_h, BRANCO)
    add_rect(s, x, y, card_w, Inches(0.55), accent)
    txt_color = BRANCO if accent in (GRAFITE, DOURADO_ESC) else GRAFITE
    add_text(s, x, y + Inches(0.12), card_w, Inches(0.4),
             nome.upper(), size=12, color=txt_color, bold=True, font=T_SANS,
             align=PP_ALIGN.CENTER)
    # Preço
    add_text(s, x, y + Inches(0.85), card_w, Inches(0.3),
             "AQUISIÇÃO", size=8.5, color=DOURADO_ESC, bold=True, font=T_SANS,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.15), card_w, Inches(0.6),
             preco, size=20, color=GRAFITE, bold=True, font=T_SERIF,
             align=PP_ALIGN.CENTER)
    # Cap Rate
    add_text(s, x, y + Inches(1.95), card_w, Inches(0.3),
             "CAP RATE", size=8.5, color=DOURADO_ESC, bold=True, font=T_SANS,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.2), card_w, Inches(0.5),
             cap, size=18, color=GRAFITE, bold=True, font=T_SERIF,
             align=PP_ALIGN.CENTER)
    # ROI
    add_text(s, x, y + Inches(2.85), card_w, Inches(0.3),
             "ROI 10 ANOS", size=8.5, color=DOURADO_ESC, bold=True, font=T_SANS,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(3.1), card_w, Inches(0.5),
             roi, size=18, color=GRAFITE, bold=True, font=T_SERIF,
             align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
         "Cenários não consideram a economia adicional do Requalifica Centro — que pode elevar o ROI em 8 a 15 pontos percentuais.",
         size=9.5, color=CINZA_M, italic=True, font=T_SANS, align=PP_ALIGN.CENTER)

page_number(s, 11, color=GRAFITE)

# ============================================================
# SLIDE 12 — SEGURANÇA JURÍDICA & PRÓXIMOS PASSOS
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_image(s, IMG / "visao varanda 2.jpg", Inches(7.5), 0, Inches(5.833), SH)
ov = add_rect(s, Inches(7.5), 0, Inches(1.0), SH, OFFWHITE)
sp = ov.fill._xPr.find(qn("a:solidFill"))
clr = sp.find(qn("a:srgbClr"))
alpha = etree.SubElement(clr, qn("a:alpha")); alpha.set("val", "50000")

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "SEGURANÇA JURÍDICA", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(8), Inches(0.8),
         "Documentação pronta para operar",
         font=T_SERIF, size=26, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

itens = [
    ("38 matrículas desmembradas", "Alienações parciais e garantias segregadas possíveis"),
    ("AVCB regular e atualizado", "Auto de Vistoria do Corpo de Bombeiros em conformidade"),
    ("Acervo técnico completo", "Projetos arquitetônicos e plantas originais disponíveis"),
    ("Tombamento como diferencial", "Elegível ao selo \"Patrimônio Histórico Requalificado\" e linhas de fomento cultural"),
    ("Pronto para garantias", "Cessão fiduciária e financiamentos estruturados com instituições de primeira linha"),
    ("Due diligence aberta", "Toda a documentação disponível mediante NDA"),
]
y = Inches(2.2)
for h, d in itens:
    add_text(s, Inches(0.7), y, Inches(0.3), Inches(0.4),
             "▎", size=16, color=DOURADO, bold=True, font=T_SANS)
    add_text(s, Inches(0.95), y, Inches(6.3), Inches(0.35),
             h, size=12, color=GRAFITE, bold=True, font=T_SANS)
    add_text(s, Inches(0.95), y + Inches(0.32), Inches(6.3), Inches(0.4),
             d, size=10, color=CINZA_M, font=T_SANS, line_spacing=1.3)
    y += Inches(0.78)

# Próximos passos
add_rect(s, Inches(0.7), Inches(6.7), Inches(6.5), Inches(0.6), GRAFITE)
add_text(s, Inches(0.9), Inches(6.78), Inches(6.5), Inches(0.4),
         "PRÓXIMOS PASSOS · visita técnica · NDA · proposta formal",
         size=11, color=DOURADO, bold=True, font=T_SANS)

page_number(s, 12, color=GRAFITE)

# ============================================================
# SLIDE 13 — CONTATO / FINALIZAÇÃO
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)

# Logo centralizado superior
add_image(s, LOGO, Inches(2.5), Inches(0.8), Inches(8.333), Inches(2.5))

# Linha dourada
gold_line(s, Inches(5.5), Inches(3.6), Inches(2.333), 2)

# Contato
add_text(s, 0, Inches(4.0), SW, Inches(0.5),
         "CONTATO", size=11, color=DOURADO_ESC, bold=True, font=T_SANS,
         align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(4.4), SW, Inches(0.8),
         "Nicollas J. Haimerl",
         font=T_SERIF, size=32, color=GRAFITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(5.4), SW, Inches(0.5),
         "nicollas.joseph@mcap.com.br",
         font=T_SANS, size=16, color=DOURADO_ESC, align=PP_ALIGN.CENTER)

# Rodapé
add_rect(s, 0, Inches(6.7), SW, Inches(0.8), GRAFITE)
add_text(s, 0, Inches(6.85), SW, Inches(0.4),
         "OPORTUNIDADE DE INVESTIMENTO IMOBILIÁRIO",
         size=11, color=DOURADO, bold=True, font=T_SANS, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(7.2), SW, Inches(0.3),
         "Praça da República, 162/166 · Centro Histórico · São Paulo/SP",
         size=9, color=OFFWHITE, font=T_SANS, align=PP_ALIGN.CENTER, italic=True)

# ============================================================
# SALVAR
# ============================================================
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"OK: {OUT}")
print(f"Slides: {len(prs.slides)}")
