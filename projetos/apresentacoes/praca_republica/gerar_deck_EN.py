"""
Generates .pptx presentation — Real Estate Investment Opportunity (ENGLISH)
Praça da República, 162/166 — Historic Center of São Paulo
13 slides, focus on SALE + Requalifica Centro Program
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).parent
IMG = Path(r"C:\Users\xandao\Desktop\predio venda")
LOGO = Path(r"C:\Users\xandao\Downloads\capa.png")
OUT = ROOT / "Real_Estate_Investment_Opportunity_Praca_Republica_EN.pptx"

# Palette
GRAFITE = RGBColor(0x1F, 0x29, 0x37)
DOURADO = RGBColor(0xC9, 0xA9, 0x61)
DOURADO_ESC = RGBColor(0x8B, 0x6F, 0x3E)
OFFWHITE = RGBColor(0xF5, 0xF1, 0xEA)
CINZA = RGBColor(0xE5, 0xE5, 0xE5)
CINZA_M = RGBColor(0x6B, 0x72, 0x80)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
VERDE = RGBColor(0x16, 0x82, 0x4F)

T_SERIF = "Georgia"
T_SANS = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh


def add_text(slide, x, y, w, h, text, *, font=T_SANS, size=14, color=GRAFITE,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_image(slide, path, x, y, w, h):
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def gold_line(slide, x, y, w, thickness_pt=2, color=DOURADO):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = Pt(thickness_pt)
    return line


def add_overlay(slide, x, y, w, h, alpha_val="60000", color=GRAFITE):
    ov = add_rect(slide, x, y, w, h, color)
    sp = ov.fill._xPr.find(qn("a:solidFill"))
    clr = sp.find(qn("a:srgbClr"))
    a = etree.SubElement(clr, qn("a:alpha"))
    a.set("val", alpha_val)
    return ov


def page_number(slide, num, total=13, color=DOURADO):
    add_text(slide, Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
             f"{num:02d} / {total:02d}", size=9, color=color, align=PP_ALIGN.RIGHT,
             font=T_SANS)


def footer_brand(slide, color=CINZA_M):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(9), Inches(0.3),
             "PRAÇA DA REPÚBLICA · 162/166 · HISTORIC CENTER · SÃO PAULO, BRAZIL",
             size=8, color=color, font=T_SANS)


# ============================================================
# SLIDE 1 — COVER
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, GRAFITE)
add_image(s, IMG / "vista do predio.jpg", Inches(5.5), 0, Inches(7.833), SH)
add_overlay(s, Inches(5.5), 0, Inches(2.5), SH, "70000")
add_rect(s, Inches(0.7), Inches(2.7), Inches(0.04), Inches(2.2), DOURADO)

add_text(s, Inches(0.9), Inches(2.5), Inches(5), Inches(0.5),
         "INVESTMENT DOSSIER", size=12, color=DOURADO, font=T_SANS, bold=True)
add_text(s, Inches(0.9), Inches(2.9), Inches(5), Inches(2.2),
         "Real Estate\nOpportunity\nin São Paulo",
         font=T_SERIF, size=44, color=OFFWHITE, bold=True, line_spacing=1.05)
add_text(s, Inches(0.9), Inches(5.3), Inches(5), Inches(0.5),
         "Praça da República · 162/166", size=16, color=OFFWHITE,
         font=T_SERIF, italic=True)
add_text(s, Inches(0.9), Inches(5.7), Inches(5), Inches(0.4),
         "Historic Center · São Paulo, Brazil", size=12, color=DOURADO, font=T_SANS)
add_text(s, Inches(0.9), Inches(6.9), Inches(6), Inches(0.3),
         "HERITAGE-LISTED · 4,000 sqm · ZEU ZONING · REQUALIFICA CENTRO PROGRAM",
         size=8.5, color=DOURADO, font=T_SANS, bold=True)

# ============================================================
# SLIDE 2 — EXECUTIVE SUMMARY
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "EXECUTIVE SUMMARY", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "A singular asset in the heart of São Paulo",
         font=T_SERIF, size=28, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

add_text(s, Inches(0.7), Inches(2.1), Inches(6.3), Inches(2.5),
         "A 4,000 sqm commercial building located on Praça da "
         "República, directly facing the subway station (Lines 3-Red "
         "and 4-Yellow). The former headquarters of the State of São "
         "Paulo's Department of Education is a heritage-listed asset "
         "with consolidated documentation and full eligibility for "
         "the Requalifica Centro Program — the most aggressive tax "
         "incentive package ever offered for downtown São Paulo.",
         size=12, color=GRAFITE, line_spacing=1.4)

def kpi_card(x, y, label, value, sub=""):
    add_rect(s, x, y, Inches(2.7), Inches(1.4), BRANCO)
    add_rect(s, x, y, Inches(0.08), Inches(1.4), DOURADO)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(2.5), Inches(0.3),
             label, size=8.5, color=CINZA_M, bold=True, font=T_SANS)
    add_text(s, x + Inches(0.2), y + Inches(0.45), Inches(2.5), Inches(0.6),
             value, size=20, color=GRAFITE, bold=True, font=T_SERIF)
    if sub:
        add_text(s, x + Inches(0.2), y + Inches(1.05), Inches(2.5), Inches(0.3),
                 sub, size=8.5, color=DOURADO_ESC, font=T_SANS, italic=True)

kpi_card(Inches(7.3), Inches(2.1), "BUILT AREA", "4,000 sqm", "13 floors + mezzanine + loft")
kpi_card(Inches(10.2), Inches(2.1), "APPRAISAL", "USD 5.6–10M*", "direct sale")
kpi_card(Inches(7.3), Inches(3.6), "TITLES", "38 deeds", "individually registered")
kpi_card(Inches(10.2), Inches(3.6), "ZONING", "ZEU", "mixed-use, high density")

add_rect(s, 0, Inches(5.4), SW, Inches(1.5), GRAFITE)
add_text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.4),
         "INVESTMENT PILLARS", size=10, color=DOURADO, bold=True, font=T_SANS)

pillars = [
    ("Location", "Facing Praça da República\nsubway 100m away"),
    ("Heritage", "Officially listed building\nhistoric value"),
    ("Structure", "38 individual deeds\nfull flexibility"),
    ("Incentives", "Requalifica Centro\nIPTU + ISS + ITBI + 25%"),
]
for i, (h, d) in enumerate(pillars):
    x = Inches(0.7 + i * 3.1)
    add_text(s, x, Inches(5.95), Inches(3), Inches(0.4),
             h.upper(), size=11, color=DOURADO, bold=True, font=T_SANS)
    add_text(s, x, Inches(6.3), Inches(3), Inches(0.8),
             d, size=10, color=OFFWHITE, font=T_SANS, line_spacing=1.3)

# nota USD
add_text(s, Inches(0.7), Inches(4.7), Inches(8), Inches(0.3),
         "* Local appraisal: BRL 31–56 million. USD figures at illustrative FX rate of 5.55.",
         size=8.5, color=CINZA_M, italic=True, font=T_SANS)

page_number(s, 2)
footer_brand(s)

# ============================================================
# SLIDE 3 — THE ASSET
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_image(s, IMG / "recepcao + elevadores.jpg", 0, 0, Inches(5.5), SH)
add_rect(s, Inches(5.5), 0, Inches(7.833), SH, OFFWHITE)
add_text(s, Inches(5.9), Inches(0.6), Inches(7), Inches(0.4),
         "THE ASSET", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(5.9), Inches(1.0), Inches(7), Inches(0.8),
         "Technical Profile", font=T_SERIF, size=30, color=GRAFITE, bold=True)
gold_line(s, Inches(5.9), Inches(1.85), Inches(2), 2.5)

profile = [
    ("Address", "Praça da República, 162/166 — Downtown, São Paulo, Brazil"),
    ("Built Area", "4,000 sqm (≈ 43,055 sqft)"),
    ("Structure", "13 floors + mezzanine + loft"),
    ("Rooms", "30 wide units, flexible layout"),
    ("Elevators", "2 operating units"),
    ("Capacity", "Up to 600 simultaneous occupants"),
    ("Zoning", "ZEU — Centrality Zone"),
    ("Titles", "38 individually registered deeds"),
    ("Fire Cert.", "AVCB — fully compliant, up to date"),
    ("History", "Former HQ — São Paulo State Dept. of Education"),
    ("Heritage", "Officially listed building"),
    ("Property Tax", "≈ BRL 13,000 / month (whole building)"),
]
y = Inches(2.15)
for label, value in profile:
    add_text(s, Inches(5.9), y, Inches(2.0), Inches(0.32),
             label.upper(), size=9, color=DOURADO_ESC, bold=True, font=T_SANS)
    add_text(s, Inches(7.95), y, Inches(5.2), Inches(0.32),
             value, size=10.5, color=GRAFITE, font=T_SANS)
    y += Inches(0.34)

page_number(s, 3, color=GRAFITE)

# ============================================================
# SLIDE 4 — LOCATION
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, GRAFITE)
add_image(s, IMG / "visao varanda.jpg", Inches(6.5), 0, Inches(6.833), SH)
add_overlay(s, Inches(6.5), 0, Inches(1.5), SH, "60000")

add_text(s, Inches(0.7), Inches(0.6), Inches(6), Inches(0.4),
         "STRATEGIC LOCATION", size=11, color=DOURADO, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(1.0), Inches(6), Inches(1.6),
         "Heart of\ndowntown SP",
         font=T_SERIF, size=36, color=OFFWHITE, bold=True, line_spacing=1.0)
gold_line(s, Inches(0.7), Inches(2.85), Inches(2), 2.5)

locs = [
    ("100 m", "from República subway station — Lines 3-Red and 4-Yellow"),
    ("Direct", "access to Ipiranga, São João, São Luís and Consolação avenues"),
    ("Nearby", "Copan Building, Mário de Andrade Library, Municipal Theater"),
    ("Inside", "the accelerated urban redevelopment perimeter"),
    ("Hub", "for commercial, legal, educational, cultural and tourism activity"),
]
y = Inches(3.2)
for n, d in locs:
    add_text(s, Inches(0.7), y, Inches(1.3), Inches(0.5),
             n, size=18, color=DOURADO, bold=True, font=T_SERIF)
    add_text(s, Inches(2.1), y + Inches(0.05), Inches(4.4), Inches(0.6),
             d, size=11, color=OFFWHITE, font=T_SANS, line_spacing=1.3)
    y += Inches(0.7)

page_number(s, 4)
footer_brand(s, OFFWHITE)

# ============================================================
# SLIDE 5 — LAYOUT & FLEXIBILITY
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "LAYOUT & FLEXIBILITY", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "30 wide rooms · 38 individually registered deeds",
         font=T_SERIF, size=26, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

add_text(s, Inches(0.7), Inches(2.1), Inches(5.5), Inches(2.2),
         "Each floor was designed to host large-scale corporate "
         "operations. The 38 independent property titles allow for "
         "partial sales, segregated collateral and flexible "
         "financial structuring — from full single-tenant sale to "
         "the assembly of a real estate fund (REIT) with multiple "
         "investors.",
         size=12, color=GRAFITE, line_spacing=1.45)

attrs = [
    "Open floor plans with widely-spaced columns",
    "Generous ceiling height",
    "Full-height windows with natural light",
    "Power distribution segmented by floor",
    "Adaptable to open space, private offices or auditorium",
]
y = Inches(4.4)
for a in attrs:
    add_text(s, Inches(0.85), y, Inches(0.3), Inches(0.3),
             "▎", size=14, color=DOURADO, bold=True, font=T_SANS)
    add_text(s, Inches(1.1), y, Inches(5.2), Inches(0.3),
             a, size=10.5, color=GRAFITE, font=T_SANS)
    y += Inches(0.36)

add_image(s, IMG / "sala grande.jpg", Inches(6.7), Inches(2.1), Inches(6.3), Inches(2.5))
add_image(s, IMG / "visao interna + externa 3 ampliada.jpg",
          Inches(6.7), Inches(4.75), Inches(6.3), Inches(2.5))

page_number(s, 5, color=GRAFITE)

# ============================================================
# SLIDE 6 — INFRASTRUCTURE
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "INFRASTRUCTURE", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "Ready to operate",
         font=T_SERIF, size=28, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

add_image(s, IMG / "elevadores.jpg", Inches(0.7), Inches(2.1), Inches(4), Inches(2.6))
add_image(s, IMG / "quadro de luz.jpg", Inches(4.85), Inches(2.1), Inches(4), Inches(2.6))
add_image(s, IMG / "corredores.jpg", Inches(9.0), Inches(2.1), Inches(3.5), Inches(2.6))

infos = [
    ("ELEVATORS", "2 operating units"),
    ("ELECTRICAL", "Power panel segmented by floor"),
    ("CIRCULATION", "Wide corridors and double staircase"),
    ("FIRE CERT.", "AVCB compliant and up to date"),
    ("HVAC", "Split + cassette systems in approved use"),
    ("TECH ARCHIVE", "Architectural plans and drawings available"),
]
y_start = Inches(5.0)
for i, (h, d) in enumerate(infos):
    col, row = i % 3, i // 3
    x = Inches(0.7 + col * 4.15)
    y = y_start + Inches(row * 1.1)
    add_rect(s, x, y, Inches(0.06), Inches(0.9), DOURADO)
    add_text(s, x + Inches(0.2), y, Inches(3.8), Inches(0.35),
             h, size=10, color=DOURADO_ESC, bold=True, font=T_SANS)
    add_text(s, x + Inches(0.2), y + Inches(0.35), Inches(3.8), Inches(0.55),
             d, size=10.5, color=GRAFITE, font=T_SANS, line_spacing=1.25)

page_number(s, 6, color=GRAFITE)

# ============================================================
# SLIDE 7 — REQUALIFICA CENTRO (COVER)
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, GRAFITE)
add_image(s, IMG / "sobreloja vista rua.jpg", Inches(6.5), 0, Inches(6.833), SH)
add_overlay(s, Inches(6.5), 0, Inches(1.0), SH, "55000")

add_text(s, Inches(0.7), Inches(0.6), Inches(7), Inches(0.4),
         "THE INVESTMENT THESIS", size=11, color=DOURADO, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(1.0), Inches(6.5), Inches(2.4),
         "Requalifica\nCentro\nProgram",
         font=T_SERIF, size=44, color=OFFWHITE, bold=True, line_spacing=0.95)
gold_line(s, Inches(0.7), Inches(4.0), Inches(2), 2.5)

add_text(s, Inches(0.7), Inches(4.25), Inches(6), Inches(2.2),
         "The asset is fully eligible under the most aggressive "
         "municipal incentive program ever launched for downtown "
         "São Paulo. The combination of tax exemptions, direct "
         "subsidy and fast-track permitting structurally transforms "
         "the investment equation.",
         size=12.5, color=OFFWHITE, line_spacing=1.45, italic=True)

add_text(s, Inches(0.7), Inches(6.5), Inches(6), Inches(0.5),
         "Legal basis: Municipal Law 17,577/2021 · Law 17,844/2022 (AIU-SCE) · "
         "Decrees 61,311/2022 and 62,465/2023 · IN SF/SUREM nº 01/2023",
         size=8.5, color=DOURADO, font=T_SANS, italic=True)

page_number(s, 7)

# ============================================================
# SLIDE 8 — TAX INCENTIVE STACK
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "REQUALIFICA CENTRO · 1/2", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "Tax incentive stack",
         font=T_SERIF, size=26, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.8), Inches(2), 2.5)

benefits = [
    {
        "tag": "IPTU",
        "title": "Property Tax",
        "body": [
            "Forgiveness of all prior tax debt",
            "100% exemption for 3 years after CCR",
            "Progressive rates from year 4 to 8",
            "Full rate only restored in year 9",
        ],
        "impact": "Savings exceeding BRL 2,000,000\nover an 8-year cycle",
    },
    {
        "tag": "ISS",
        "title": "Service Tax",
        "body": [
            "Rate reduced from 5% to 2%",
            "Applies to construction, engineering,",
            "architecture and refurbishment",
            "linked to the retrofit",
        ],
        "impact": "Savings of BRL 300,000\non a BRL 10M retrofit",
    },
    {
        "tag": "ITBI",
        "title": "Transfer Tax",
        "body": [
            "Full exemption on transfers",
            "destined for requalification",
            "Drastically reduces acquisition",
            "and structuring costs",
        ],
        "impact": "100% exemption\non acquisition",
    },
]
card_w = Inches(4.0)
card_h = Inches(4.5)
gap = Inches(0.25)
total_w = card_w * 3 + gap * 2
start_x = (SW - total_w) / 2
y = Inches(2.1)

for i, b in enumerate(benefits):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, y, card_w, card_h, BRANCO)
    add_rect(s, x, y, card_w, Inches(0.5), DOURADO)
    add_text(s, x, y + Inches(0.08), card_w, Inches(0.4),
             b["tag"], size=14, color=BRANCO, bold=True, font=T_SERIF, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + Inches(0.7), card_w - Inches(0.6), Inches(0.5),
             b["title"], size=14, color=GRAFITE, bold=True, font=T_SERIF)
    yy = y + Inches(1.3)
    for line in b["body"]:
        add_text(s, x + Inches(0.4), yy, Inches(0.2), Inches(0.3),
                 "•", size=12, color=DOURADO, bold=True, font=T_SANS)
        add_text(s, x + Inches(0.6), yy, card_w - Inches(0.8), Inches(0.3),
                 line, size=10.5, color=GRAFITE, font=T_SANS)
        yy += Inches(0.35)
    add_rect(s, x + Inches(0.3), y + card_h - Inches(1.1),
             card_w - Inches(0.6), Inches(0.85), OFFWHITE)
    add_text(s, x + Inches(0.3), y + card_h - Inches(1.05),
             card_w - Inches(0.6), Inches(0.3),
             "DIRECT IMPACT", size=8, color=DOURADO_ESC, bold=True, font=T_SANS,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + card_h - Inches(0.75),
             card_w - Inches(0.6), Inches(0.55),
             b["impact"], size=11, color=GRAFITE, bold=True, font=T_SERIF,
             align=PP_ALIGN.CENTER, line_spacing=1.15)

add_rect(s, 0, Inches(6.8), SW, Inches(0.7), GRAFITE)
add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
         "Current property tax ≈ BRL 13,000/month · 3-year exemption = BRL 468,000 saved on IPTU alone",
         size=12, color=DOURADO, bold=True, font=T_SANS, align=PP_ALIGN.CENTER)

page_number(s, 8, color=GRAFITE)

# ============================================================
# SLIDE 9 — DIRECT SUBSIDY + FAST-TRACK PERMITTING
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "REQUALIFICA CENTRO · 2/2", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "Direct subsidy + fast track",
         font=T_SERIF, size=26, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.8), Inches(2), 2.5)

add_rect(s, Inches(0.7), Inches(2.15), Inches(6.3), Inches(4.0), GRAFITE)
add_text(s, Inches(1.0), Inches(2.3), Inches(6), Inches(0.4),
         "DIRECT GOVERNMENT SUBSIDY", size=10, color=DOURADO, bold=True, font=T_SANS)
add_text(s, Inches(1.0), Inches(2.7), Inches(6), Inches(1.3),
         "Up to 25%\nnon-refundable",
         font=T_SERIF, size=34, color=OFFWHITE, bold=True, line_spacing=1.0)
add_text(s, Inches(1.0), Inches(4.5), Inches(6), Inches(1.5),
         "Direct municipal grant over construction costs, awarded "
         "through public calls issued by the Municipal Department "
         "of Urbanism and Licensing (SMUL).",
         size=11, color=OFFWHITE, font=T_SANS, line_spacing=1.4)
add_rect(s, Inches(1.0), Inches(5.5), Inches(5.7), Inches(0.5), DOURADO)
add_text(s, Inches(1.0), Inches(5.55), Inches(5.7), Inches(0.4),
         "On a BRL 10M retrofit → up to BRL 2,500,000 non-refundable",
         size=11.5, color=GRAFITE, bold=True, font=T_SANS, align=PP_ALIGN.CENTER)

add_text(s, Inches(7.3), Inches(2.3), Inches(6), Inches(0.4),
         "FAST-TRACK PERMITTING", size=10, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(7.3), Inches(2.65), Inches(6), Inches(0.6),
         "Summary rite · 60 business days",
         font=T_SERIF, size=20, color=GRAFITE, bold=True)

items = [
    "Project approval in up to 60 business days",
    "Municipal licensing fees waived for 5 years",
    "Inspection fees waived during the same period",
    "Priority status in public-sector contracts",
]
y = Inches(3.5)
for it in items:
    add_text(s, Inches(7.3), y, Inches(0.3), Inches(0.3),
             "▎", size=14, color=DOURADO, bold=True, font=T_SANS)
    add_text(s, Inches(7.55), y, Inches(5.6), Inches(0.5),
             it, size=11, color=GRAFITE, font=T_SANS, line_spacing=1.3)
    y += Inches(0.5)

add_rect(s, Inches(7.3), Inches(5.7), Inches(5.7), Inches(1.2), VERDE)
add_text(s, Inches(7.5), Inches(5.8), Inches(5.4), Inches(0.4),
         "✓ ELIGIBILITY CONFIRMED", size=10, color=BRANCO, bold=True, font=T_SANS)
add_text(s, Inches(7.5), Inches(6.15), Inches(5.4), Inches(0.7),
         "Pre-1992 building + ZEU Zone\nArticle 37 of Law 17,844/2022",
         size=10.5, color=BRANCO, font=T_SANS, line_spacing=1.3)

page_number(s, 9, color=GRAFITE)

# ============================================================
# SLIDE 10 — USE CASES
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "USE CASES", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "Multiple operating models",
         font=T_SERIF, size=28, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

uses = [
    ("Corporate Headquarters", "For companies seeking representativeness, centrality and a prestige address in São Paulo."),
    ("REIT / Real Estate Fund", "Ideal asset for income portfolio composition — 38 deeds enable phased acquisitions."),
    ("Education Operation", "Structure compatible with universities, training centers, technical schools and open courses."),
    ("Mixed-Use Retrofit", "Conversion to residential, hospitality or mixed use with enhanced municipal incentives."),
    ("Family Office", "Patrimonial positioning in a historic asset with cultural status and high fractional liquidity."),
    ("Institutional HQ", "Natural vocation for government bodies, cultural institutions, foundations or NGOs."),
]
for i, (h, d) in enumerate(uses):
    col, row = i % 2, i // 2
    x = Inches(0.7 + col * 6.3)
    y = Inches(2.15 + row * 1.55)
    add_rect(s, x, y, Inches(6), Inches(1.35), BRANCO)
    add_rect(s, x, y, Inches(0.08), Inches(1.35), DOURADO)
    add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(5.6), Inches(0.45),
             h, size=14, color=GRAFITE, bold=True, font=T_SERIF)
    add_text(s, x + Inches(0.25), y + Inches(0.6), Inches(5.6), Inches(0.75),
             d, size=10.5, color=CINZA_M, font=T_SANS, line_spacing=1.35)

page_number(s, 10, color=GRAFITE)

# ============================================================
# SLIDE 11 — RETURN SCENARIOS
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "PROJECTED RETURNS", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
         "Investment return scenarios",
         font=T_SERIF, size=28, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

add_text(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.5),
         "Simulation after acquisition and retrofit · based on local m² appreciation",
         size=11, color=CINZA_M, italic=True, font=T_SANS)

scenarios = [
    ("Conservative", "BRL 31 M",   "5.8%", "58%",  CINZA),
    ("Moderate",     "BRL 43.5 M", "6.1%", "61%",  DOURADO),
    ("Realistic",    "BRL 43.5 M", "7.4%", "74%",  DOURADO_ESC),
    ("Optimistic",   "BRL 56 M",   "6.9%", "69%",  GRAFITE),
]
card_w = Inches(2.9)
card_h = Inches(3.6)
gap = Inches(0.2)
start_x = (SW - (card_w * 4 + gap * 3)) / 2
y = Inches(2.8)

for i, (name, price, cap, roi, accent) in enumerate(scenarios):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, y, card_w, card_h, BRANCO)
    add_rect(s, x, y, card_w, Inches(0.55), accent)
    txt_color = BRANCO if accent in (GRAFITE, DOURADO_ESC) else GRAFITE
    add_text(s, x, y + Inches(0.12), card_w, Inches(0.4),
             name.upper(), size=12, color=txt_color, bold=True, font=T_SANS,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.85), card_w, Inches(0.3),
             "ACQUISITION", size=8.5, color=DOURADO_ESC, bold=True, font=T_SANS,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.15), card_w, Inches(0.6),
             price, size=20, color=GRAFITE, bold=True, font=T_SERIF,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.95), card_w, Inches(0.3),
             "CAP RATE", size=8.5, color=DOURADO_ESC, bold=True, font=T_SANS,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.2), card_w, Inches(0.5),
             cap, size=18, color=GRAFITE, bold=True, font=T_SERIF,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.85), card_w, Inches(0.3),
             "10-YEAR ROI", size=8.5, color=DOURADO_ESC, bold=True, font=T_SANS,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(3.1), card_w, Inches(0.5),
             roi, size=18, color=GRAFITE, bold=True, font=T_SERIF,
             align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
         "Scenarios do not include additional Requalifica Centro savings — which may add 8 to 15 percentage points to ROI.",
         size=9.5, color=CINZA_M, italic=True, font=T_SANS, align=PP_ALIGN.CENTER)

page_number(s, 11, color=GRAFITE)

# ============================================================
# SLIDE 12 — LEGAL SECURITY & NEXT STEPS
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)
add_image(s, IMG / "visao varanda 2.jpg", Inches(7.5), 0, Inches(5.833), SH)
add_overlay(s, Inches(7.5), 0, Inches(1.0), SH, "50000", OFFWHITE)

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
         "LEGAL SECURITY", size=11, color=DOURADO_ESC, bold=True, font=T_SANS)
add_text(s, Inches(0.7), Inches(0.9), Inches(8), Inches(0.8),
         "Documentation ready to operate",
         font=T_SERIF, size=26, color=GRAFITE, bold=True)
gold_line(s, Inches(0.7), Inches(1.85), Inches(2), 2.5)

items = [
    ("38 individual deeds", "Enables partial sales and segregated collateral structures"),
    ("AVCB compliant", "Fire department certification active and up to date"),
    ("Complete tech archive", "Original architectural plans and drawings available"),
    ("Heritage as upside", "Eligible for the \"Requalified Heritage\" seal and cultural funding lines"),
    ("Collateral-ready", "Compatible with fiduciary assignment and structured financing"),
    ("Open due diligence", "Full documentation made available under NDA"),
]
y = Inches(2.2)
for h, d in items:
    add_text(s, Inches(0.7), y, Inches(0.3), Inches(0.4),
             "▎", size=16, color=DOURADO, bold=True, font=T_SANS)
    add_text(s, Inches(0.95), y, Inches(6.3), Inches(0.35),
             h, size=12, color=GRAFITE, bold=True, font=T_SANS)
    add_text(s, Inches(0.95), y + Inches(0.32), Inches(6.3), Inches(0.4),
             d, size=10, color=CINZA_M, font=T_SANS, line_spacing=1.3)
    y += Inches(0.78)

add_rect(s, Inches(0.7), Inches(6.7), Inches(6.5), Inches(0.6), GRAFITE)
add_text(s, Inches(0.9), Inches(6.78), Inches(6.5), Inches(0.4),
         "NEXT STEPS · site visit · NDA · formal offer",
         size=11, color=DOURADO, bold=True, font=T_SANS)

page_number(s, 12, color=GRAFITE)

# ============================================================
# SLIDE 13 — CONTACT
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, OFFWHITE)

add_image(s, LOGO, Inches(2.5), Inches(0.8), Inches(8.333), Inches(2.5))
gold_line(s, Inches(5.5), Inches(3.6), Inches(2.333), 2)

add_text(s, 0, Inches(4.0), SW, Inches(0.5),
         "CONTACT", size=11, color=DOURADO_ESC, bold=True, font=T_SANS,
         align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(4.4), SW, Inches(0.8),
         "Nicollas J. Haimerl",
         font=T_SERIF, size=32, color=GRAFITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(5.4), SW, Inches(0.5),
         "nicollas.joseph@mcap.com.br",
         font=T_SANS, size=16, color=DOURADO_ESC, align=PP_ALIGN.CENTER)

add_rect(s, 0, Inches(6.7), SW, Inches(0.8), GRAFITE)
add_text(s, 0, Inches(6.85), SW, Inches(0.4),
         "REAL ESTATE INVESTMENT OPPORTUNITY",
         size=11, color=DOURADO, bold=True, font=T_SANS, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(7.2), SW, Inches(0.3),
         "Praça da República, 162/166 · Historic Center · São Paulo, Brazil",
         size=9, color=OFFWHITE, font=T_SANS, align=PP_ALIGN.CENTER, italic=True)

prs.save(str(OUT))
print(f"OK: {OUT}")
print(f"Slides: {len(prs.slides)}")
